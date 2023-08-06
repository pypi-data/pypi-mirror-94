# =======================================================
#
# =======================================================
import pandas as pd
import numpy as np
import scipy.stats as ss
import matplotlib.pyplot as plt
import seaborn as sns

from statsmodels.tsa.ar_model import AR
from statsmodels.tsa.arima_model import ARMA
from statsmodels.tsa.holtwinters import SimpleExpSmoothing
from statsmodels.tsa.holtwinters import ExponentialSmoothing
from statsmodels.tsa.statespace.sarimax import SARIMAX
from statsmodels.tsa.arima_model import ARIMA

from cryptography.fernet import Fernet
import nltk

from sklearn.ensemble import RandomForestClassifier
from sklearn.ensemble import RandomForestRegressor
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from sklearn.model_selection import train_test_split
from sklearn.model_selection import cross_val_score
from sklearn.metrics import confusion_matrix
from sklearn.metrics import classification_report
from sklearn.metrics import accuracy_score
from sklearn.metrics import roc_curve, auc
from sklearn.metrics import precision_recall_curve
from sklearn.metrics import average_precision_score
from sklearn.metrics import plot_confusion_matrix

from imblearn.ensemble import BalancedRandomForestClassifier
from sklearn.datasets import make_classification

from pptx.util import Pt
from pptx.util import Inches

# =======================================================
#
# =======================================================
def get_equal_groups(data_df,group_count,variable_name):
    """ Text """

    groups_dict = {}
    total_size = 0
    group_count_now = 1

    data_small_df = data_df[[variable_name]]
    groups_df = data_small_df.groupby([variable_name]).size()
    groups_df = groups_df.reset_index()
    avg_group_size_fl= (groups_df.sum()/group_count)[0]

    for i in range(0,len(groups_df.index)):
        group_int = (((groups_df.iloc[[i],[0]]).values)[0])[0]
        size_int  = (((groups_df.iloc[[i],[1]]).values)[0])[0]
        total_size = total_size + size_int

        groups_dict.update({group_int:group_count_now})

        if total_size > avg_group_size_fl:
            total_size = 0
            group_count_now = group_count_now + 1
    col_name = variable_name+'_gp'
    dict_ps = pd.Series(groups_dict)
    dict_df = dict_ps.to_frame()
    dict_df.columns = [col_name]

    temp_df = dict_df.reset_index(inplace=False)
    min_df = temp_df.groupby([col_name]).min().astype(str)
    max_df = temp_df.groupby([col_name]).max().astype(str)

    full_gp_name_df = pd.merge(min_df, max_df, left_index=True, right_index=True)

    full_gp_name_df[(col_name+'_rg')] = (full_gp_name_df['index_x']) + "-To-" + (full_gp_name_df['index_y'])
    return_df = pd.merge(full_gp_name_df, temp_df, left_on=col_name, right_on=col_name)
    return_df = return_df[[col_name,(col_name+'_rg'),'index']]
    return_df = return_df.set_index('index')

    return return_df

# =======================================================
#
# =======================================================
def get_col_summary(data_df,variable,variable_name,label_variable_name):
    """ Text """
    new1_df = data_df.groupby([variable_name]).count()
    new2_df = data_df.groupby([variable_name,label_variable_name]).count()

    t_df = (new2_df[variable] / new1_df[variable])
    t_df = t_df.reset_index()
    pivot_df = t_df.pivot(index=variable_name, columns=label_variable_name, values=variable)

    t1 = new1_df[label_variable_name].reset_index()
    t1['TotalSize'] = t1[label_variable_name] / t1[label_variable_name].sum()
    return_df = pd.merge(t1, pivot_df, left_on=variable_name, right_on=variable_name)

    return return_df

# =======================================================
#
# =======================================================
def analyze_df(data_df, label_name,group_cnt_int):
    """ Text """
    final_df = pd.DataFrame()
    counter = 0
    columns_list = data_df.columns.tolist()
    columns_list.remove(label_name)
    for entry in columns_list:
        columnd_dtype_st = str(data_df[entry].dtypes)
        if columnd_dtype_st == 'int64' or columnd_dtype_st =='float64':
            groups_df = get_equal_groups(data_df,group_cnt_int,entry)
            data_df = pd.merge(data_df, groups_df, left_on=entry, right_index=True)
            new_col_name = entry+'_gp_rg'
            col_pivot = get_col_summary(data_df,entry,new_col_name,label_name)
            col_pivot['col_name'] = new_col_name
            col_pivot.rename(columns={ col_pivot.columns[0]: "var_name" }, inplace = True)
            if counter == 0:
                final_df = col_pivot
                counter = 1
            else:
                final_df = final_df.append(col_pivot)

    return final_df,data_df

# =======================================================
#
# =======================================================
def cramers( x_in , y_in):
    """ Text """
    confusion_m = pd.crosstab(x_in,y_in)
    chi2 = ss.chi2_contingency(confusion_m)[0]

    n_val = confusion_m.sum().sum()
    n_val_1 = n_val-1

    phi2 = chi2/n_val

    r,k = confusion_m.shape

    ph2c = max(0, phi2-((k-1)*(r-1))/n_val_1)

    rcorr = r-((r-1)**2)/n_val_1
    kcorr = k-((k-1)**2)/n_val_1

    kc_1 = (kcorr-1)
    rc_1 = (rcorr-1)

    ret_val = np.sqrt(ph2c/min(kc_1,rc_1))

    return ret_val

# =======================================================
#
# =======================================================
def analyze_association(df,Label_column,verbose):
    """ Text """
    for col in df.columns:
        if len(df[col].unique()) == 1:
            df.drop(col,inplace=True,axis=1)

    final = []

    col_num_list = list(df.select_dtypes('number').columns)
    col_obj_list = list(df.select_dtypes('object').columns)
    col_cat_list = list(df.select_dtypes('category').columns)
    col_cat_list = col_cat_list + col_obj_list

    if verbose == 1:
        print('======================NumToNum======================')
    for col1 in col_num_list:
        for col2 in col_num_list:
            if verbose == 1:
                print(col2+"|"+col1)
            x = df[col1]
            y = df[col2]
            final.append(['cramers_v',col1,col2,cramers(x, y)])
            final.append(['cramers_v',col2,col1,cramers(y, x)])

    if verbose == 1:
        print('======================CatToNum======================')
    for col1 in col_cat_list:
        for col2 in col_num_list:
            if verbose == 1:
                print(col2+"|"+col1)

            df.fillna(value={col2: 999},inplace=True)
            #FIX HERE TO FILL CAT VALUES
            #df.fillna(value={col1: 999},inplace=True)

            x = df[col1]
            y = df[col2]
            final.append(['cramers_v',col1,col2,cramers(x, y)])
            final.append(['cramers_v',col2,col1,cramers(y, x)])

    if verbose == 1:
        print('======================CatToCat======================')
    for col1 in col_cat_list:
        for col2 in col_cat_list:
            if verbose == 1:
                print(col2+"|"+col1)
            x = df[col1]
            y = df[col2]
            final.append(['cramers_v',col1,col2,cramers(x, y)])
            final.append(['cramers_v',col2,col1,cramers(y, x)])

    df = pd.DataFrame.from_records(final,columns=['algo','col1','col2','val'])

    df = df[df['algo'] == 'cramers_v']
    df = df.drop(['algo'], axis=1)
    df = df.drop_duplicates()

    DATA_RECORDS = df.pivot(index='col2', columns='col1', values='val')

    if verbose == 1:
        print (DATA_RECORDS[Label_column].sort_values(ascending=False)[:30], '\n')

    return DATA_RECORDS

# =======================================================
#
# =======================================================
def get_important_features(df,Label_Column):
    """ Text """
    # ================================================================================================
    # ================================================================================================
    nunique = df.apply(pd.Series.nunique)
    cols_to_drop = nunique[nunique <= 1].index
    df = df.drop(cols_to_drop, axis=1)
	# ================================================================================================
    # ================================================================================================
    df = df.apply(lambda x:x.fillna(x.value_counts().index[0]))
    features_df = df.drop(Label_Column,axis=1,inplace=False)
    labels_df = df[Label_Column]
    onehot_df = pd.get_dummies(features_df)
    features_np = onehot_df.values.tolist()
    labels_np = labels_df.values.tolist()
    # ================================================================================================
    # ================================================================================================

    if str(df[Label_Column].dtypes) != 'float64':
        clf = RandomForestClassifier(n_estimators=10,n_jobs=-1,verbose=1)
        RandomForestClassifier
        print("using RandomForestClassifier")
    else :
        clf = RandomForestRegressor(n_estimators=10,n_jobs=-1,verbose=1)
        RandomForestRegressor
        print("using RandomForestRegressor")

    clf = clf.fit(features_np,labels_np)

    importances = clf.feature_importances_
    importances_df = pd.DataFrame(importances)

    column_names_df  = pd.DataFrame(onehot_df.columns.values.tolist())

    importances_df_final = pd.merge(column_names_df,importances_df, left_index=True,right_index=True)
    importances_df_final.columns = ["Feature_Name", "Importance"]
    importances_df_final = importances_df_final.sort_values(by=['Importance'], ascending=False)

    return importances_df_final

# =======================================================
#
# =======================================================
def analyze_correlation(df,Label_column):
    """ Text """
    df_local = df

    label_type = str(df_local[Label_column].dtypes)
    if (label_type != 'int64') and (label_type != 'float64') and (label_type != 'int32'):
        labels = df_local[Label_column].unique().tolist()
        mapping = dict( zip(labels,range(len(labels))) )
        df_local = df_local.replace({Label_column: mapping},inplace=False)

    numeric_features = df_local.select_dtypes(include=[np.number])
    corr = numeric_features.corr()

    return corr

# =======================================================
#
# =======================================================
def predict_1(data_list,month_range):
    """ Text """
    model = AR(data_list)
    model_fit = model.fit()
    y = model_fit.predict(len(data_list), (len(data_list)+month_range))

    return y

# =======================================================
#
# =======================================================
def predict_2(data_list,month_range):
    """ Text """
    model = ARMA(data_list,order=(0,1))
    model_fit = model.fit(disp=False)
    y = model_fit.predict(len(data_list), (len(data_list)+month_range))

    return y

# =======================================================
#
# =======================================================
def predict_6(data_list,month_range):
    """ Text """
    model = ARIMA(data_list)
    model_fit = model.fit(disp=False)
    y = model_fit.predict(len(data_list), (len(data_list)+month_range))

    return y

# =======================================================
#
# =======================================================
def predict_3(data_list,month_range):
    """ Text """
    model = ExponentialSmoothing(data_list)
    model_fit = model.fit()
    y = model_fit.predict(len(data_list), (len(data_list)+month_range))

    return y

# =======================================================
#
# =======================================================
def predict_4(data_list,month_range):
    """ Text """
    model = SARIMAX(data_list,order=(1,1,1), seasonal_order=(1, 1, 1, 1))
    model_fit = model.fit()
    y = model_fit.predict(len(data_list), (len(data_list)+month_range))

    return y

# =======================================================
#
# =======================================================
def predict_5(data_list,month_range):
    """ Text """
    model = SimpleExpSmoothing(data_list)
    model_fit = model.fit()
    y = model_fit.predict(len(data_list), (len(data_list)+month_range))

    return y

# =======================================================
#
# =======================================================
def append_forecast(df,mer_list,churn_list,month_range):
    """ Text """
    row_cnt = len(df.index) +1
    prev = 0
    for i in range(row_cnt,(month_range+1)):
        prev = i -1
        mer = mer_list[i-row_cnt]
        churn = churn_list[i-row_cnt]

        prev_deact = df.loc[prev]['deact']
        prev_cnt = df.loc[prev]['cnt']
        prev_cnt = prev_cnt - (prev_cnt * churn) /100

        ABSdeact = (prev_cnt * churn)/100
        prev_deact = prev_deact + ABSdeact

        churn2 = 1- (prev_deact / df.loc[1]['cnt'])

        df_row = pd.DataFrame([[i,prev_cnt,mer,churn,ABSdeact,prev_deact,churn2]],columns=['month_n','cnt','MER','ABSchurn','ABSdeact','deact','churn'])
        df_row = df_row.set_index('month_n')
        df = df.append(df_row,sort=True)

    return df

# =======================================================
#
# =======================================================
def calculate_LTV(df,month_range):
    """ Text """
    mer_sum = df["MER"].mean()
    final_churn = df["ABSchurn"].mean()/100
    LTV = (mer_sum/final_churn)

    return LTV,final_churn,mer_sum

# =======================================================
#
# =======================================================
def calculate_ltv(df):
    """ Text """
    LTV_List = []
    filename_list = []
    counter = 0
    month_range = 48

    ga_date = sorted(df.start_date.unique())
    for date in ga_date:
        df_GA = df[df['start_date'] == date].copy()

        #TOTAL
        df_TOTAL = df_GA.groupby(['month_n']).sum()
        df_TOTAL = df_TOTAL.drop(['svc_rate'],axis=1)

        df_GA['REV'] = df_GA['cnt'] * df_GA['svc_rate']
        df_GA = df_GA.groupby(['servicescancellationdate','month_n']).sum()
        df_GA['MER'] = df_GA['REV'] / df_GA['cnt']
        df_GA = df_GA.drop(['svc_rate','REV'],axis=1)
        df_GA = df_GA.loc['A']

        df_GA['churn'] = df_GA['cnt'] / df_TOTAL['cnt']
        df_GA['deact'] = df_TOTAL['cnt'] - df_GA['cnt']
        df_GA['ABSdeact'] = df_GA.diff()['deact']
        df_GA['ABSchurn'] = (df_GA['ABSdeact'] / df_GA['cnt']) * 100
        df_GA = df_GA.fillna(0)

        mer_list = predict_5(df_GA['MER'].tolist(),month_range)
        churn_list = predict_2(df_GA['ABSchurn'].tolist(),month_range)
        final_df = append_forecast(df_GA,mer_list,churn_list,month_range)
        LTV,final_churn,mer_sum = calculate_LTV(final_df,month_range)
        size = df_GA.loc[1]['cnt']
        active = final_df.loc[(month_range)]['cnt']
        LTV_List.append([LTV,date,(final_churn*100),mer_sum,size,active])

        final_df['SP_DATE'] = date

        fig, ax1 = plt.subplots(figsize=(6, 3), dpi=150)
        x = final_df.index
        y1 = final_df.MER
        y2 = final_df.ABSchurn
        y4 = (final_df.deact/final_df['cnt'][1]) * 100

        if counter == 0:
            df_run_total = final_df
            Acct_sum_df = (pd.DataFrame([y1,y2,y4])).T
            Acct_sum_df['date'] = str(date)
            Acct_sum_df['ltv'] = str(round((LTV),2))
            Acct_sum_df['total'] = str(final_df['cnt'][1])
            counter = 1
        else:
            df_run_total=  df_run_total.append(final_df)
            Acct_tmp_df = (pd.DataFrame([y1,y2,y4])).T
            Acct_tmp_df['date'] = date
            Acct_tmp_df['ltv'] = str(round((LTV),2))
            Acct_tmp_df['total'] = str(final_df['cnt'][1])
            Acct_sum_df = Acct_sum_df.append(Acct_tmp_df)

        ax2 = ax1.twinx()
        ax1.plot(x,y1,'g-', label='MER $')
        ax1.plot(x,y2,'b-', label='Churn %')
        ax2.plot(x,y4,'y-', label='Deact (y2)')
        ax1.set_xticks(np.arange(1, 50, 4))
        leg=ax2.legend( frameon=False)
        leg=ax1.legend(frameon=False)
        ax1.set_xlabel('Self-Pay Month')
        ax1.set_ylabel('MER $ & CHURN %')
        ax2.set_ylabel('Deact %')
        title_str = "Self Pay Date=" + str(date) +"|LTV=$"+str(int(LTV))+"|Churn="+str(round((final_churn*100),2))+"%" +"|N="+str(int(final_df['cnt'][1]))
        ax1.axvline(x=len(df_GA.index), color='r', linestyle='--',alpha=.5)
        plt.title(title_str)
        plt.show()
        filename = str(len(df_GA.index)) +"-"+ date + 'LTV.png'
        fig.savefig(filename,transparent=False,  bbox_inches='tight',pad_inches=0.15,dpi=200)
        filename_list.append(filename)

    final_results_sum = pd.DataFrame.from_records(LTV_List,columns=['LTV','SP_Date','Churn','MER','Size','active'])

    return Acct_sum_df, final_results_sum,filename_list

# =======================================================
#
# =======================================================
def get_heatmap(df,filename='heat_map.png',scale=1,font_size=12,dec_space='0.2f',auto_scale=0,dpi=300,figsizex=14,figsizey=12):
    """ Text """
    sns.set(font_scale=scale)
    fig = plt.figure(figsize=(figsizex, figsizey),frameon =False)
    fig = sns.heatmap(df,annot=True,cbar=False,linewidths=0.4,annot_kws={"size": font_size},fmt=dec_space, square=False,cmap="Blues")
    plt.show()
    fig = fig.get_figure()
    fig.savefig(filename,transparent=False,  bbox_inches='tight',pad_inches=0.05,dpi=dpi)
    plt.close()

# =======================================================
#
# =======================================================
def get_bar(df,filename,x_name,y_name,figsizex=15,figsizey=15):
    """ Text """
    sns.set(style="whitegrid")
    ax = plt.figure(figsize=(figsizex, figsizey),frameon =False)
    ax = sns.barplot(df[x_name], df[y_name], palette="GnBu_d")
    plt.show()
    ax.figure.savefig(filename,transparent=False,  bbox_inches='tight',pad_inches=0.05,dpi=300)
    plt.close()

# =======================================================
#
# =======================================================
def get_anal_col(df,date_col,label_col):
    """ Text """
    col_list = list(df.columns.values)
    col_list.remove(date_col)
    col_list.remove(label_col)

    return col_list

# =======================================================
#
# =======================================================
def pivot_by_col(df,date_col,label_col,pivot_col):
    """ Text """
    df_new = df.copy()

    df_new = df_new[[date_col,label_col,pivot_col]]
    df_new.fillna(value=np.nan, inplace=True)

    df_new[pivot_col] = df_new[pivot_col].astype(str)

    total_size = df_new.groupby([pivot_col,date_col]).size()
    total_size_conv = df_new[df_new[label_col] == 1].groupby([pivot_col,date_col]).size()
    conv_df = (total_size_conv/total_size) * 100

    total_size = df_new.groupby([date_col]).size()
    total_size_conv = df_new.groupby([pivot_col,date_col]).size()
    distro_df = (total_size_conv/total_size) * 100

    return pd.DataFrame(conv_df),pd.DataFrame(distro_df),pd.DataFrame(total_size_conv)

# =======================================================
#
# =======================================================
def ACA_create_graphs(df,date,label):
    """ Text """
    col_list = get_anal_col(df,date,label)
    Chart_list = []
    for entry in col_list:
        print(entry)
        conv_df,distro_df,total_size_df = pivot_by_col(df,date,label,entry)

        conv_df.reset_index(inplace = True)
        conv_df.rename(columns={ conv_df.columns[2]: "val" }, inplace = True)
        pivot_df = conv_df.pivot(index=entry, columns=date, values='val')
        get_heatmap(pivot_df,entry+"conv_df.png",0.4,6,'0.2f',0,300,8,4)

        distro_df.reset_index(inplace = True)
        distro_df.rename(columns={ distro_df.columns[2]: "val" }, inplace = True)
        pivot_df = distro_df.pivot(index=entry, columns=date, values='val')
        get_heatmap(pivot_df,entry+"dist_df.png",0.4,6,'0.2f',0,300,8,4)

        total_size_df.reset_index(inplace = True)
        total_size_df.rename(columns={ total_size_df.columns[2]: "val" }, inplace = True)
        pivot_df = total_size_df.pivot(index=entry, columns=date, values='val')
        get_heatmap(pivot_df,entry+"size_df.png",0.4,6,'0.0f',0,300,8,4)

        Chart_list.append([entry+"conv_df.png",'Conversion'])
        Chart_list.append([entry+"dist_df.png",'Distribution'])
        Chart_list.append([entry+"size_df.png",'Size'])

    return Chart_list

# =======================================================
#
# =======================================================
def add_chart_slide(prs,image_name,title):
    """ Text """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = Inches(0.05)
    top =  Inches(0.03)
    width = Inches(9.8)
    height = Inches(7.5)

    try:
        slide.shapes.add_picture(image_name,left,top,width)
    except:
        print("ERROR")

    left = Inches(5)
    top =  Inches(6.5)
    width = Inches(1)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(title)
    p.font.size = Pt(10)
# =======================================================
#
# =======================================================
def add_table_slide(prs,df,title):
    """ Text """
    df = df.round(2)
    print(df)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    left = Inches(0.1)
    top =  Inches(0.6)
    width = Inches(9.5)
    height = Inches(3)

    table = shapes.add_table(17,5,left,top,width,height).table

    # write column headings
    table.cell(0, 0).text = 'LTV $'
    table.cell(0, 1).text = 'SP DATE'
    table.cell(0, 2).text = 'CHURN %'
    table.cell(0, 3).text = 'MER $'
    table.cell(0, 4).text = 'SIZE #'

    len_df = len(df)
    if len_df > 17:
        len_df = 17

    # write body cells
    for i in range(1,len_df):
        table.cell(i,0).text = str((((df.iloc[[i],[0]]).values)[0])[0])
        table.cell(i,1).text = str((((df.iloc[[i],[1]]).values)[0])[0])
        table.cell(i,2).text = str((((df.iloc[[i],[2]]).values)[0])[0] )
        table.cell(i,3).text = str((((df.iloc[[i],[3]]).values)[0])[0] )
        table.cell(i,4).text = str((((df.iloc[[i],[4]]).values)[0])[0] )

    left = Inches(0.05)
    top =  Inches(0.05)
    width = Inches(1)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(title)
    p.font.size = Pt(10)

# =======================================================
#
# =======================================================
def create_title(prs,title_txt):
    """ Text """
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = title_txt

# =======================================================
#
# =======================================================
def save_presentation(prs,filename):
    """ Text """
    filename = filename +  '.pptx'
    prs.save(filename)

# =======================================================
#
# =======================================================
def decrypt_df(transcript,key):
    """ Text """
    cipher_suite = Fernet(key)
    tBytes = bytes(transcript,'UTF-8')
    decrypted = cipher_suite.decrypt(tBytes)

    return(decrypted.decode('UTF-8'))

# =======================================================
#
# =======================================================
def encrypt_df(transcript,key):
    """ Text """
    cipher_suite = Fernet(key)
    tBytes = bytes(transcript,'UTF-8')
    encrypted = cipher_suite.encrypt(tBytes)

    return(encrypted.decode('UTF-8'))

# =======================================================
#
# =======================================================
def nlp_text_process(note,proc_type):
    """ Text """
    words = note.split(' ')

    if proc_type == 'stem':
        porter = nltk.PorterStemmer()
        word_list = [porter.stem(t) for t in words]

    if proc_type == 'lem':
        WNlemma = nltk.WordNetLemmatizer()
        word_list = [WNlemma.lemmatize(t) for t in words]

    sent = " ".join(word_list)

    return sent
# =======================================================
#
# =======================================================
def nlp_split(note,end_str,start_str):
    """ Text """
    start_pos = 0
    end_pos = 0
    final_note = ""
    len_note = len(note)

    while end_pos <= len_note:
        start = note.find(start_str,start_pos)
        end = note.find(end_str,end_pos)
        note_part = note[start:end]
        final_note = final_note +"|"+ note_part
        if end < end_pos:
            end = len_note
        start_pos = end
        end_pos = end+1

    return final_note
# =======================================================
#
# =======================================================
def nlp_cluster(data_df,cluster_column,clusters,prediction_col,max_df_val,min_df_val,max_iter_val,n_init_val,cluster_algo,ngram_range_val):
    """ Text """

    vectorizer = TfidfVectorizer(stop_words='english',use_idf = True, lowercase  = True, max_df = max_df_val, min_df = min_df_val, ngram_range = ngram_range_val)
    X = vectorizer.fit_transform(data_df[cluster_column].values)

    model = KMeans(n_clusters=clusters, init='k-means++', max_iter=max_iter_val, n_init=n_init_val, algorithm='auto',n_jobs = -1,verbose=True)
    model.fit(X)
    labels = model.labels_.tolist()

    data_df[prediction_col] = labels

    return data_df,model,X,vectorizer
# =======================================================
#
# =======================================================
def make_confusionMatrix_graph(confusionMatrix):
    """ Text """
    plt.matshow(confusionMatrix)
    plt.title('Confusion matrix')
    plt.colorbar()
    plt.ylabel('True label')
    plt.xlabel('Predicted label')
    fn = "conf_matrix_graph.png"
    plt.savefig(fn)
    plt.close()

    return fn

# =======================================================
#
# =======================================================
def make_precision_recall_graph(labels_test,prediction):
    """ Text """
    precision, recall, thresholds = precision_recall_curve(labels_test, prediction)
    average_precision = average_precision_score(labels_test, prediction)
    plt.clf()
    plt.plot(recall, precision, color='navy', label='Precision-Recall curve')
    plt.xlabel('Recall')
    plt.ylabel('Precision')
    plt.title('Precision-Recall curve|area ='+str(average_precision))
    fn = "precision_recall.png"
    plt.savefig(fn)
    plt.close()
    return fn

# =======================================================
#
# =======================================================
def make_ROC_graph(labels_test,prediction):
    """ Text """
    false_positive_rate, recall, thresholds = roc_curve(labels_test,prediction)
    roc_auc = auc(false_positive_rate, recall)
    plt.title('Receiver Operating Characteristic')
    plt.plot(false_positive_rate, recall, 'b', label='AUC = %0.2f' %roc_auc)
    plt.legend(loc='lower right')
    plt.plot([0, 1], [0, 1], 'r--')
    plt.xlim([0.0, 1.0])
    plt.ylim([0.0, 1.0])
    plt.ylabel('Recall')
    plt.xlabel('Fall-out')
    fn = "roc_graph.png"
    plt.savefig(fn)
    plt.close()

    return fn

# =======================================================
#
# =======================================================
def run_classification(df,Label_Column,logger,is_forecast_flag,run_name,estimators,max_depth_i,run_inbalance):
    """ Text """
    logger.info('========================================================')
    logger.info('========================================================')
    logger.info(run_name)
    logger.info('run_classification')

    # ================================================================================================
    # ================================================================================================

    logger.info('Remove Dups')
    nunique = df.apply(pd.Series.nunique)
    cols_to_drop = nunique[nunique <= 1].index
    df = df.drop(cols_to_drop, axis=1)

    # ================================================================================================
    # ================================================================================================

    logger.info('One Hot Encoding')
    df = df.apply(lambda x:x.fillna(x.value_counts().index[0]))

    is_forecast_index = df[df[is_forecast_flag] == 1].index
    is_train_index = df[df[is_forecast_flag] == 0].index

    features_df = df.drop(Label_Column,axis=1,inplace=False)
    features_df = features_df.drop(is_forecast_flag,axis=1,inplace=False)

    labels_df = df[[Label_Column]]
    onehot_df = pd.get_dummies(features_df)

    onehot_df_train   = onehot_df.loc[is_train_index]
    labels_df_train   = labels_df.loc[is_train_index]

    onehot_df_forecast = onehot_df.loc[is_forecast_index]

    features_np = onehot_df_train.values.tolist()
    features_np_forecast = onehot_df_forecast.values.tolist()
    labels_np = labels_df_train.values.ravel()

    # ================================================================================================
    # ================================================================================================
    features_train, features_test, labels_train, labels_test = train_test_split(features_np,labels_np, random_state=0,test_size=0.2)
    # ================================================================================================
    # ================================================================================================
    if run_inbalance == 1:
        clf = BalancedRandomForestClassifier(n_estimators=estimators,n_jobs=-1,verbose=1,max_depth=max_depth_i)
    if run_inbalance == 0:
        clf = RandomForestClassifier(n_estimators=estimators,n_jobs=-1,verbose=1,max_depth=max_depth_i)
    RandomForestClassifier
    logger.info("using RandomForestClassifier")

    # ================================================================================================
    # ================================================================================================

    logger.info('Train Model')
    clf = clf.fit(features_train,labels_train)
    scores = cross_val_score(clf, features_train, labels_train)
    logger.info("cross_val_score for train data"+str(scores))

    # ================================================================================================
    # ================================================================================================

    prediction = clf.predict(features_test)
    scores = cross_val_score(clf, features_test, labels_test)
    logger.info("cross_val_score for test data"+str(scores))
    confusionMatrix = confusion_matrix(labels_test,prediction)
    tn,fp,fn,tp = confusion_matrix(labels_test,prediction).ravel()

    # ================================================================================================
    # ================================================================================================

    g1 = make_confusionMatrix_graph(confusionMatrix)
    g2 = make_ROC_graph(labels_test,prediction)
    g3 = make_precision_recall_graph(labels_test,prediction)

    target_names = ['0','1']

    logger.info("Classification Report_TEST \n"+str(classification_report(labels_test, prediction, target_names=target_names)))
    logger.info("Accuracy Score_TEST= \n" +str(accuracy_score(labels_test, prediction)))
    logger.info("Confusion Matrix_TEST \n"+str(confusionMatrix))
    logger.info("tn ="+str(tn))
    logger.info("fp ="+str(fp))
    logger.info("fn ="+str(fn))
    logger.info("tp ="+str(tp))

    # ================================================================================================
    # ================================================================================================

    importances = clf.feature_importances_
    importances_df = pd.DataFrame(importances)
    column_names_df  = pd.DataFrame(onehot_df.columns.values.tolist())
    importances_df_final = pd.merge(column_names_df,importances_df, left_index=True,right_index=True)
    importances_df_final.columns = ["Feature_Name", "Importance"]
    importances_df_final = importances_df_final.sort_values(by=['Importance'], ascending=False)
    imp_list = [importances_df_final.columns.values.tolist()] + importances_df_final.values.tolist()
    importances_df_final.to_csv('imp_features.csv', index=False,)
    importances_df_final.to_excel('imp_features.xlsx', index=False,)

    # ================================================================================================
    # ================================================================================================
    logger.info('Predict')

    prediction = clf.predict(features_np_forecast)
    probability = clf.predict_proba(features_np_forecast)

    probability_df = pd.DataFrame(probability)
    probability_df.columns = ['Prob_0','Prob_1']

    probability_df['orig_index']= is_forecast_index
    probability_df['prediction']= prediction

    probability_df = probability_df.set_index('orig_index')

    # ================================================================================================
    # ================================================================================================
    return importances_df_final,probability_df,g1,g2,g3,run_name,clf,column_names_df,target_names,confusionMatrix

# =======================================================
#
# =======================================================
